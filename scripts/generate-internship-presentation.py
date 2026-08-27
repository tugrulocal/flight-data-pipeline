#!/usr/bin/env python3
"""Flight Data Pipeline staj sunumunu düzenlenebilir PowerPoint olarak üretir."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


OUT = Path("deliverables/flight-data-pipeline-staj-sunumu.pptx")
P = Presentation()
P.slide_width = Inches(13.333)
P.slide_height = Inches(7.5)

NAVY = "102A43"
BLUE = "1976D2"
CYAN = "00ACC1"
GREEN = "2E7D32"
ORANGE = "EF6C00"
RED = "C62828"
PURPLE = "6A1B9A"
GREY = "5B6770"
LIGHT = "F5F8FB"
DARK = "172B4D"
WHITE = "FFFFFF"


def rgb(value):
    return RGBColor.from_string(value)


def rect(slide, x, y, w, h, color=WHITE, radius=False, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.color.rgb = rgb(line or color)
    return shape


def text(slide, value, x, y, w, h, size=18, color=DARK, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font="Aptos", margin=0.08):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def title(slide, heading, subtitle=None, number=None):
    rect(slide, 0, 0, 13.333, 0.18, BLUE)
    text(slide, heading, 0.62, 0.36, 11.5, 0.52, 26, NAVY, True)
    if subtitle:
        text(slide, subtitle, 0.64, 0.93, 11.5, 0.34, 11, GREY)
    if number:
        text(slide, f"{number:02d}", 12.15, 0.37, 0.55, 0.36, 12, BLUE, True, PP_ALIGN.RIGHT)
    rect(slide, 0.62, 7.1, 12.1, 0.01, "D9E2EC")
    text(slide, "Flight Data Pipeline • Staj Projesi", 0.62, 7.17, 5.3, 0.18, 8, GREY)


def bullets(slide, items, x, y, w, h, size=17, color=DARK, bullet_color=BLUE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(11)
    return box


def card(slide, heading, body, x, y, w, h, accent=BLUE):
    rect(slide, x, y, w, h, LIGHT, True, "D9E2EC")
    rect(slide, x, y, 0.08, h, accent)
    text(slide, heading, x + 0.25, y + 0.18, w - 0.4, 0.35, 15, NAVY, True)
    text(slide, body, x + 0.25, y + 0.63, w - 0.42, h - 0.76, 12, GREY)


def arrow(slide, x1, y1, x2, y2, color=BLUE, width=2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def node(slide, label, x, y, w, h, color=BLUE, note=None):
    rect(slide, x, y, w, h, color, True, color)
    text(slide, label, x + 0.08, y + 0.12, w - 0.16, h * 0.48, 14, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    if note:
        text(slide, note, x + 0.08, y + h * 0.56, w - 0.16, h * 0.32, 9, WHITE, False, PP_ALIGN.CENTER)


def agenda_slide():
    s = P.slides.add_slide(P.slide_layouts[6])
    title(s, "Sunum akışı", "15–20 dakika anlatım + 3–4 dakika canlı demo", 2)
    agenda = [
        ("1", "Problem ve hedef"), ("2", "Canlı veri mimarisi"),
        ("3", "Güvenilirlik ve hata yönetimi"), ("4", "Kullanıcı arayüzü"),
        ("5", "DevOps, Kubernetes ve izleme"), ("6", "Spark / Iceberg analitik katmanı"),
        ("7", "Test, release ve sonraki adımlar"),
    ]
    for i, (n, label) in enumerate(agenda):
        col, row = i % 2, i // 2
        x, y = 0.95 + col * 6.1, 1.55 + row * 1.25
        rect(s, x, y, 5.55, 0.9, LIGHT, True, "D9E2EC")
        rect(s, x + 0.18, y + 0.17, 0.54, 0.54, BLUE if i < 4 else PURPLE, True)
        text(s, n, x + 0.18, y + 0.24, 0.54, 0.25, 13, WHITE, True, PP_ALIGN.CENTER)
        text(s, label, x + 0.9, y + 0.25, 4.3, 0.3, 18, NAVY, True)
    return s


def add_content():
    # 1 Cover
    s = P.slides.add_slide(P.slide_layouts[6])
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    rect(s, 0, 0, 13.333, 0.18, CYAN)
    text(s, "Flight Data Pipeline", 0.75, 1.25, 8.8, 0.72, 34, WHITE, True)
    text(s, "Canlı uçuş verisinden güvenilir veri hattına,\nweb arayüzüne ve analitik katmana", 0.78, 2.15, 8.4, 1.0, 21, "D9EAF7")
    text(s, "Staj Projesi Sunumu", 0.8, 3.55, 4.5, 0.35, 16, CYAN, True)
    text(s, "OpenSky • Kafka • MongoDB • FastAPI • React • Docker • Kubernetes • Spark • Iceberg", 0.8, 4.1, 10.6, 0.3, 13, WHITE)
    for i, (lab, col) in enumerate([("API", CYAN), ("STREAM", BLUE), ("DB", GREEN), ("WEB", ORANGE), ("ANALYTICS", PURPLE)]):
        node(s, lab, 0.8 + i * 2.25, 5.45, 1.85, 0.72, col)
        if i:
            arrow(s, 0.8 + i * 2.25 - 0.3, 5.81, 0.8 + i * 2.25, 5.81, "8FA7BF")

    agenda_slide()

    # 3 Problem
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "Başlangıç problemi", "Canlı dış veriyi kullanıcıya doğrudan bağlamak yeterli değildir.", 3)
    card(s, "Dış kaynak değişken", "OpenSky API gecikebilir, kota uygulayabilir veya geçici olarak yanıt vermeyebilir.", 0.75, 1.55, 3.8, 2.15, ORANGE)
    card(s, "Canlı ekran hassas", "Kullanıcı haritayı açtığında güncel bir durum görmeli; bağlantı koparsa toparlanmalıdır.", 4.77, 1.55, 3.8, 2.15, BLUE)
    card(s, "Veri kaybı riski", "Bir servis çökerken mesajı kaybetmeden, tekrar işlese bile çift kayıt üretmeden devam etmeliyiz.", 8.78, 1.55, 3.8, 2.15, RED)
    text(s, "Çözüm yaklaşımı", 0.78, 4.45, 2.4, 0.35, 20, NAVY, True)
    bullets(s, ["Canlı veri ile kalıcı veriyi ayırmak", "Servisleri mesaj kuyruğu ile gevşek bağlamak", "Hataları görünür ve tekrar işlenebilir kılmak", "Kullanıcıya REST + WebSocket ile kontrollü veri sunmak"], 0.85, 4.95, 10.8, 1.65, 17)

    # 4 Goals
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "Proje hedefleri", "Tek hedef ekran yapmak değil; uçtan uca veri mühendisliği pratiği kazanmaktı.", 4)
    card(s, "Ürün hedefi", "Canlı uçak listesi, konum haritası, seçili uçağın rotası ve geçmiş sorgusu.", 0.8, 1.45, 5.65, 1.55, BLUE)
    card(s, "Öğrenme hedefi", "Docker, Kafka, MongoDB, FastAPI, React, test, release ve operasyon kavramlarını birlikte uygulamak.", 6.85, 1.45, 5.65, 1.55, GREEN)
    card(s, "Kalite hedefi", "Veri kaybına dayanıklılık, idempotent yazım, hata izolasyonu, doğrulanabilir kurulum ve geri alma yolu.", 0.8, 3.45, 5.65, 1.55, ORANGE)
    card(s, "Analitik hedef", "Ham olaylardan Spark ile Bronze / Silver / Gold katmanları ve Iceberg snapshot yönetimi oluşturmak.", 6.85, 3.45, 5.65, 1.55, PURPLE)
    text(s, "Başarı ölçütü: Kodun çalışması + gözlemlenebilirlik + test kanıtı + güvenli işletim.", 0.86, 5.7, 11.4, 0.45, 18, NAVY, True, PP_ALIGN.CENTER)

    # 5 architecture
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "Canlı veri hattı: büyük resim", "Veri kaynağından kullanıcı ekranına kadar sorumluluklar ayrılmıştır.", 5)
    node(s, "OpenSky", 0.5, 2.55, 1.45, 0.95, CYAN, "REST API")
    node(s, "Producer", 2.35, 2.55, 1.55, 0.95, BLUE, "Python")
    node(s, "Kafka", 4.35, 2.55, 1.55, 0.95, ORANGE, "raw topic")
    node(s, "MongoDB\nConsumer", 6.35, 1.5, 1.75, 0.95, GREEN, "kalıcı yazım")
    node(s, "FastAPI", 6.35, 3.65, 1.75, 0.95, PURPLE, "REST + WS")
    node(s, "MongoDB", 8.6, 1.5, 1.55, 0.95, GREEN, "raw + live")
    node(s, "React", 10.65, 3.65, 1.55, 0.95, BLUE, "harita")
    arrow(s, 1.95, 3.03, 2.35, 3.03); arrow(s, 3.9, 3.03, 4.35, 3.03)
    arrow(s, 5.9, 2.85, 6.35, 2.0); arrow(s, 5.9, 3.2, 6.35, 4.13)
    arrow(s, 8.1, 1.98, 8.6, 1.98); arrow(s, 8.1, 4.13, 10.65, 4.13)
    arrow(s, 9.38, 2.45, 8.1, 4.0, GREEN)
    text(s, "İki ayrı tüketici grubu aynı mesajları bağımsız okur; biri veritabanına yazar, diğeri canlı yayın yapar.", 0.78, 5.55, 11.8, 0.55, 17, NAVY, True, PP_ALIGN.CENTER)

    # 6 Producer
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "1. Veri üretimi: OpenSky Producer", "Dış API yanıtını ortak bir olay sözleşmesine dönüştüren Python servisi.", 6)
    bullets(s, ["OpenSky REST API’den uçak konumlarını periyodik olarak alır.", "Global veya Türkiye + yakın çevre bounding-box modunda çalışabilir.", "Her olaya schema_version ve benzersiz event_id ekler.", "Kimlik bilgileri ve çalışma ayarları kodda değil, ortam değişkenlerinde tutulur.", "Rate-limit bilgisine göre bekler; anonim kullanımda güvenli minimum sorgu aralığı uygular."], 0.75, 1.45, 6.4, 4.5, 16)
    card(s, "Neden event_id?", "Mesaj yeniden işlenirse aynı gerçek olayı yeniden tanıyabilmek için kalıcı kimlik gerekir.", 7.65, 1.55, 4.6, 1.45, BLUE)
    card(s, "Hata halinde", "API geçici yanıt vermezse producer kontrollü yeniden dener; Kafka zaten almış mesajlar kaybolmaz.", 7.65, 3.3, 4.6, 1.45, ORANGE)
    card(s, "Doğrulama", "Producer loglarında poll sayısı / gönderilen mesaj; Kafka topic’inde yeni kayıtlar izlenir.", 7.65, 5.05, 4.6, 1.1, GREEN)

    # 7 kafka
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "2. Kafka: veri hattının omurgası", "Kafka, üretici ile tüketicileri birbirinden bağımsız çalıştıran kalıcı mesaj katmanıdır.", 7)
    node(s, "Producer", 0.8, 2.6, 1.8, 1.0, BLUE, "üretir")
    node(s, "aircraft.positions\n.raw.v1", 4.1, 2.45, 2.4, 1.3, ORANGE, "Kafka topic")
    node(s, "Mongo writer\ngroup", 8.0, 1.45, 1.85, 1.0, GREEN, "tüm mesajlar")
    node(s, "Realtime gateway\ngroup", 8.0, 3.85, 1.85, 1.0, PURPLE, "tüm mesajlar")
    arrow(s, 2.6, 3.1, 4.1, 3.1); arrow(s, 6.5, 2.88, 8.0, 1.96); arrow(s, 6.5, 3.32, 8.0, 4.34)
    card(s, "Önemli ayrım", "Aynı consumer group olsalardı mesajları paylaşırlar ve iki işlevin de her mesajı görmesi garanti olmazdı.", 0.8, 5.35, 10.7, 0.95, RED)
    text(s, "Ayrı group → aynı olayın iki farklı iş amacı için bağımsız okunması", 1.1, 6.45, 10.8, 0.3, 16, NAVY, True, PP_ALIGN.CENTER)

    # 8 Mongo
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "3. MongoDB: geçmiş ve canlı durum ayrımı", "Tek koleksiyonla hem arşiv hem ekran performansı hedeflenmedi.", 8)
    card(s, "raw_positions", "Her Kafka olayının tarihçesi\n• _id = event_id\n• rota / analiz / hata inceleme\n• index: icao24 + observed_at\n• 48 saat TTL", 0.9, 1.6, 5.25, 3.2, GREEN)
    card(s, "live_positions", "Her uçak için son durum\n• _id = icao24\n• frontend ilk snapshot kaynağı\n• hızlı canlı liste sorgusu\n• 7 gün TTL", 7.15, 1.6, 5.25, 3.2, BLUE)
    arrow(s, 6.15, 3.15, 7.15, 3.15, "8FA7BF")
    text(s, "Ham olay = analiz ve geçmiş • Son durum = canlı kullanıcı deneyimi", 1.1, 5.55, 11.0, 0.45, 19, NAVY, True, PP_ALIGN.CENTER)

    # 9 delivery
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "4. Teslim garantisi: neden veri kaybolmuyor?", "Tasarım: at-least-once teslim + idempotent MongoDB yazımı.", 9)
    steps = [("1", "Kafka mesajı alınır", BLUE), ("2", "raw_positions upsert", GREEN), ("3", "live_positions upsert", GREEN), ("4", "Offset commit", ORANGE)]
    for i, (n, lab, col) in enumerate(steps):
        x = 0.65 + i * 3.12
        rect(s, x, 2.05, 2.55, 1.25, LIGHT, True, "D9E2EC")
        rect(s, x + 0.14, 2.31, 0.55, 0.55, col, True)
        text(s, n, x + 0.14, 2.43, 0.55, 0.22, 13, WHITE, True, PP_ALIGN.CENTER)
        text(s, lab, x + 0.85, 2.4, 1.5, 0.35, 14, NAVY, True, PP_ALIGN.CENTER)
        if i < 3: arrow(s, x + 2.55, 2.68, x + 3.12, 2.68)
    card(s, "Çökme örneği", "MongoDB yazımından sonra fakat offset commit’ten önce süreç çökerse Kafka mesajı yeniden gönderir.", 0.95, 4.1, 5.45, 1.35, ORANGE)
    card(s, "Tekrarın güvenliği", "event_id ile yapılan upsert, aynı Kafka mesajının ikinci kez ayrı raw belge yaratmasını önler.", 6.85, 4.1, 5.45, 1.35, GREEN)
    text(s, "Bu, uçtan uca exactly-once değildir; fakat pratikte veri kaybı ve duplicate riskini yönetir.", 1.0, 6.05, 11.3, 0.35, 16, RED, True, PP_ALIGN.CENTER)

    # 10 DLQ retry
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "5. Bozuk mesaj ve geçici hata yönetimi", "Ana akışın durmaması, hatanın görünmez olması anlamına gelmez.", 10)
    node(s, "Geçerli mesaj", 0.8, 2.25, 1.8, 0.9, GREEN)
    node(s, "MongoDB yazımı", 3.4, 2.25, 1.9, 0.9, BLUE)
    node(s, "Başarılı\ncommit", 6.2, 1.2, 1.7, 0.9, GREEN)
    node(s, "Geçici hata\nretry", 6.2, 3.3, 1.7, 0.9, ORANGE)
    node(s, "Kalıcı bozuk\nDLQ", 9.0, 3.3, 1.7, 0.9, RED)
    arrow(s, 2.6, 2.7, 3.4, 2.7); arrow(s, 5.3, 2.45, 6.2, 1.65); arrow(s, 5.3, 2.95, 6.2, 3.75); arrow(s, 7.9, 3.75, 9.0, 3.75)
    bullets(s, ["Geçici MongoDB hatalarında sınırlı exponential backoff ile retry yapılır.", "Kalıcı biçimde bozuk mesajlar Base64 kaynak zarfıyla DLQ topic’ine taşınır.", "DLQ teslimi doğrulanmadan ana offset ilerletilmez.", "DLQ operasyonel olarak izlenmelidir; hata saklanmaz, ayrıştırılır."], 0.9, 5.15, 11.2, 1.25, 14)

    # 11 backend
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "6. FastAPI: frontend için güvenli geçit", "Tarayıcı Kafka’ya veya MongoDB’ye doğrudan bağlanmaz.", 11)
    node(s, "MongoDB", 0.85, 2.15, 1.8, 0.9, GREEN, "snapshot + history")
    node(s, "Kafka", 0.85, 4.2, 1.8, 0.9, ORANGE, "realtime group")
    node(s, "FastAPI", 4.5, 3.15, 2.1, 1.15, PURPLE, "REST + WebSocket")
    node(s, "React / Nginx", 9.15, 3.15, 2.1, 1.15, BLUE, "tarayıcı")
    arrow(s, 2.65, 2.6, 4.5, 3.58); arrow(s, 2.65, 4.65, 4.5, 3.85); arrow(s, 6.6, 3.72, 9.15, 3.72)
    card(s, "REST", "/health, /api/aircraft, tek uçak, geçmiş, istatistik. İlk açılıştaki güvenilir durum kaynağıdır.", 0.8, 5.55, 5.35, 0.95, BLUE)
    card(s, "WebSocket", "/ws/aircraft. Yeni uçak konumlarını anlık iter; bağlantı sonrası REST ile yeniden eşitlenir.", 6.85, 5.55, 5.35, 0.95, PURPLE)

    # 12 Frontend
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "7. React frontend: canlı, anlaşılır ve dayanıklı", "Kullanıcının gördüğü ekran, gerçek zamanlı veri hattının son adımıdır.", 12)
    card(s, "Canlı harita", "MapLibre varsayılan motor. Uçak sembolleri uçuş yönüne göre döner; irtifa renkleri yüksek kontrastlıdır.", 0.75, 1.45, 3.8, 2.05, BLUE)
    card(s, "Geçmiş rota", "Uçağa tıklanınca raw_positions üzerinden geçmiş noktalar okunur; rota irtifa aralığına göre renklidir.", 4.78, 1.45, 3.8, 2.05, GREEN)
    card(s, "Uyumluluk", "WebGL ya da yükleme hatasında Leaflet otomatik fallback olur; kullanıcıya yeniden deneme seçeneği sunulur.", 8.8, 1.45, 3.8, 2.05, ORANGE)
    bullets(s, ["Sadece son 20 dakika içinde gözlenen uçaklar canlı ekranda gösterilir.", "WebSocket mesajları requestAnimationFrame ile toplu uygulanır; her mesaj ayrı render yaratmaz.", "Nginx aynı-origin proxy ile REST ve WebSocket’i frontend üzerinden güvenli biçimde iletir."], 0.9, 4.2, 11.2, 1.7, 16)

    # 13 Docker
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "8. Docker Compose: taşınabilir yerel sistem", "Her bileşen ayrı container; bütün sistem tek Compose tanımıyla ayağa kalkar.", 13)
    components = [("kafka-volume-init", RED), ("Kafka", ORANGE), ("topic-init", ORANGE), ("MongoDB", GREEN), ("producer", BLUE), ("consumer", GREEN), ("backend", PURPLE), ("frontend", BLUE)]
    for i, (lab, col) in enumerate(components):
        x, y = 0.8 + (i % 4) * 3.05, 1.55 + (i // 4) * 1.35
        node(s, lab, x, y, 2.45, 0.8, col)
    card(s, "Kalıcı veri", "Kafka ve MongoDB named volume kullanır. Normal docker compose down veriyi korur; down -v veri siler.", 0.8, 4.6, 5.4, 1.25, GREEN)
    card(s, "Ağ yüzeyi", "Release Compose yalnız frontend’i 127.0.0.1:5175 üzerinde host’a açar. Backend/Kafka/Mongo iç ağda kalır.", 6.85, 4.6, 5.4, 1.25, BLUE)
    text(s, "Temiz Kafka volume’larda izin sorununu çözmek için kafka-volume-init, broker başlatılmadan yalnız volume sahibini düzeltir.", 0.85, 6.25, 11.5, 0.38, 14, NAVY, True, PP_ALIGN.CENTER)

    # 14 Kubernetes
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "9. Kubernetes öğrenme yolu ve GitOps", "Yerelde çalışan yapı, Kubernetes kaynaklarına aşamalı olarak taşındı.", 14)
    lessons = [("01", "MongoDB + temel servisler"), ("02", "Kafka + topic init"), ("03", "Producer + Consumer"), ("04", "Backend + Frontend + Ingress"), ("05", "Tam sistem + izleme")]
    for i, (num, lab) in enumerate(lessons):
        x = 0.65 + i * 2.55
        rect(s, x, 2.0, 2.05, 1.55, LIGHT, True, "D9E2EC")
        rect(s, x + 0.73, 1.72, 0.6, 0.6, PURPLE, True)
        text(s, num, x + 0.73, 1.85, 0.6, 0.22, 12, WHITE, True, PP_ALIGN.CENTER)
        text(s, lab, x + 0.18, 2.48, 1.7, 0.55, 12, NAVY, True, PP_ALIGN.CENTER)
        if i < 4: arrow(s, x + 2.05, 2.78, x + 2.55, 2.78, "8FA7BF")
    card(s, "Deployment", "Producer, consumer, backend ve frontend çökerse Kubernetes yeni pod başlatır.", 0.85, 4.55, 3.6, 1.15, BLUE)
    card(s, "StatefulSet + PVC", "Kafka ve MongoDB yeniden başlasa bile kalıcı diskleri korunur.", 4.85, 4.55, 3.6, 1.15, GREEN)
    card(s, "Argo CD", "Manifestler Git üzerinden izlenir ve cluster’a GitOps yaklaşımıyla uygulanır.", 8.85, 4.55, 3.6, 1.15, PURPLE)

    # 15 Monitoring
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "10. İzleme: sistemin sağlığını görünür kılmak", "Çalışan sistem ile gözlemlenebilir sistem aynı şey değildir.", 15)
    node(s, "Producer", 0.85, 2.55, 1.65, 0.85, BLUE, "metrik")
    node(s, "Backend", 3.2, 2.55, 1.65, 0.85, PURPLE, "metrik")
    node(s, "Prometheus", 5.6, 2.55, 1.9, 0.85, ORANGE, "toplar")
    node(s, "Grafana", 8.3, 2.55, 1.65, 0.85, GREEN, "görselleştirir")
    arrow(s, 2.5, 2.98, 5.6, 2.8); arrow(s, 4.85, 3.12, 5.6, 3.12); arrow(s, 7.5, 2.98, 8.3, 2.98)
    bullets(s, ["Prometheus ve Grafana ayrı monitoring namespace’inde yönetilir.", "Pod hazır olma, CPU/bellek/ağ ve uygulama metrikleri izlenebilir.", "Grafana dashboard değişiklikleri de kod ve Argo CD ile takip edilir.", "/health ile readiness aynı anlamda değildir: uygulama ayakta olabilir, dış bağımlılıkları eksik olabilir."], 0.9, 4.55, 11.3, 1.55, 16)

    # 16 analytics flow
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "11. Analitik katman: canlı hattan bağımsız Spark", "Analiz işini canlı uçuş deneyimini etkilemeden ayrı bir batch yolunda ele aldık.", 16)
    node(s, "MongoDB\nraw_positions", 0.65, 2.55, 1.75, 1.0, GREEN, "salt-okunur")
    node(s, "JSONL.gz\nexport", 3.3, 2.55, 1.65, 1.0, BLUE, "Git dışı")
    node(s, "Spark local[2]", 5.75, 2.55, 1.85, 1.0, ORANGE, "Docker")
    node(s, "Bronze / Silver\n/ Rejected", 8.45, 1.45, 1.9, 1.0, PURPLE, "detay")
    node(s, "Gold\nözetler", 8.45, 3.75, 1.9, 1.0, GREEN, "rapor")
    arrow(s, 2.4, 3.05, 3.3, 3.05); arrow(s, 4.95, 3.05, 5.75, 3.05); arrow(s, 7.6, 2.75, 8.45, 1.95); arrow(s, 7.6, 3.35, 8.45, 4.25)
    text(s, "Önemli sınır: Spark bu aşamada Kafka’yı tüketmez ve MongoDB’ye yazmaz; canlı yol aynen çalışmaya devam eder.", 0.95, 5.7, 11.4, 0.45, 16, RED, True, PP_ALIGN.CENTER)

    # 17 medalion
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "12. Bronze / Silver / Gold: veriyi adım adım olgunlaştırmak", "Her katmanın farklı bir sorumluluğu vardır.", 17)
    card(s, "Bronze — ham ama tiplenmiş", "Dışa aktarılan bütün okunabilir olaylar. Kalite filtresi uygulanmaz; veri kaynağına en yakın katmandır.", 0.75, 1.75, 3.7, 3.35, ORANGE)
    card(s, "Silver — analiz için güvenilir", "Kimlik, zaman, koordinat ve hız kurallarını geçen olaylar. Hatalı / eski gözlemler rejected tablosuna ayrılır.", 4.82, 1.75, 3.7, 3.35, "7B1FA2")
    card(s, "Gold — karar için özet", "Saatlik trafik, aktivite ve veri kalitesi gibi dashboard/rapor dostu sonuçlar yalnız Silver’dan üretilir.", 8.9, 1.75, 3.7, 3.35, GREEN)
    arrow(s, 4.45, 3.42, 4.82, 3.42); arrow(s, 8.52, 3.42, 8.9, 3.42)
    text(s, "Örnek kalite kuralları: negatif veya 400 m/s üzeri hız; observed_at ile ingested_at arasında 20 dakikayı aşan gecikme.", 0.9, 5.8, 11.6, 0.45, 15, NAVY, True, PP_ALIGN.CENTER)

    # 18 Iceberg
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "13. Apache Iceberg: tablo sürümü ve güvenli yenileme", "Iceberg, Parquet dosyalarının üzerinde şema, partition ve snapshot yönetimi sağlar.", 18)
    card(s, "Tablo soyutlaması", "Veri dosyaları Parquet olarak kalır; metadata klasörü şemayı, snapshot’ları ve hangi dosyaların tabloda olduğunu yönetir.", 0.85, 1.5, 5.35, 1.55, PURPLE)
    card(s, "Snapshot / time-travel", "Başarılı append yeni snapshot oluşturur. Güncel tablo normal SQL ile; geçmiş hal VERSION AS OF ile sorgulanabilir.", 6.85, 1.5, 5.35, 1.55, BLUE)
    card(s, "Kontrollü batch refresh", "Önce dry-run; sonra event_id anti-join ile yalnız yeni olaylar Bronze’a append edilir. İki kez ekleme engellenir.", 0.85, 3.65, 5.35, 1.55, GREEN)
    card(s, "Kaynak sınırı", "Tek Docker container içinde Spark local[2], en fazla 2 CPU ve 3 GiB bellek kullanır; Kubernetes kaynaklarını korur.", 6.85, 3.65, 5.35, 1.55, ORANGE)
    text(s, "Amaç: ham veriye dokunmadan, tekrar çalıştırılabilir ve izlenebilir analiz yenilemeleri yapmak.", 0.85, 6.2, 11.6, 0.4, 17, NAVY, True, PP_ALIGN.CENTER)

    # 19 lifecycle
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "14. Test yaklaşımı: her katmanı kanıtlamak", "Bir değişiklik ancak gözlemlenebilir kanıtla tamamlanmış sayılır.", 19)
    rows = [
        ("Kod", "Python syntax, unit test, frontend test/build", BLUE),
        ("Entegrasyon", "İzole Kafka + Mongo: retry, DLQ, lag, TTL, retention, index", GREEN),
        ("Uçtan uca", "REST snapshot, WebSocket, canlı harita, history rotası", PURPLE),
        ("Operasyon", "restart, backup / boş hedef restore, log ve consumer lag", ORANGE),
        ("Release", "çok mimari image, SBOM, provenance, Cosign, offline paket", RED),
    ]
    for i, (h, b, col) in enumerate(rows):
        y = 1.35 + i * 1.05
        rect(s, 0.9, y, 1.7, 0.72, col, True)
        text(s, h, 0.98, y + 0.2, 1.54, 0.23, 14, WHITE, True, PP_ALIGN.CENTER)
        rect(s, 2.85, y, 9.3, 0.72, LIGHT, True, "D9E2EC")
        text(s, b, 3.1, y + 0.18, 8.85, 0.28, 14, NAVY)

    # 20 Release
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "15. Release ve dağıtım disiplini", "Geliştirme ortamında çalışan şey, temiz hedefte de doğrulanmalıdır.", 20)
    bullets(s, ["AMD64 ve ARM64 için image’lar ayrı üretilir; registry’de çok mimarili manifest hedeflenir.", "SBOM (bileşen listesi), provenance (üretim kökeni) ve keyless Cosign imzası release güven zincirinin parçasıdır.", "Offline paket dış/iç SHA-256 doğrulaması, sürümlü Compose ve setup/backup/restore scriptleri içerir.", "Setup scripti mimari, Docker/Compose, bellek, port, disk ve image bütünlüğünü denetler.", "Rollback: APP_VERSION önceki sürüme alınabilir; veri şeması için backup/restore prosedürü ayrıca korunur."], 0.75, 1.4, 7.0, 4.9, 16)
    card(s, "RC durumu", "v1.0.0-rc.2 temiz Kafka volume izin hatası nedeniyle kabul edilmedi; etiket değiştirilmez.", 8.25, 1.6, 3.9, 1.35, RED)
    card(s, "RC.3 kapısı", "kafka-volume-init ile temiz volume testleri; registry/offline, manifest/SBOM/provenance/Cosign doğrulaması.", 8.25, 3.25, 3.9, 1.45, ORANGE)
    card(s, "v1.0.0", "Bütün kabul satırları geçer ve kullanıcı onay verirse aynı uygulama commit’i etiketlenir.", 8.25, 5.05, 3.9, 1.15, GREEN)

    # 21 Demo
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "16. Canlı demo planı", "Demo kısa tutulur: ürün değerini ve hattın sağlığını birlikte gösterir.", 21)
    demo = [("1", "Servisler", "docker compose ps ile container health durumunu göster."), ("2", "Kullanıcı ekranı", "Haritayı aç; canlı uçak listesini ve filtrelemeyi göster."), ("3", "Geçmiş rota", "Bir uçağı seç; geçmiş noktaların rota olarak çizildiğini göster."), ("4", "Backend kanıtı", "/health ve /api/aircraft endpoint’lerini göster."), ("5", "Veri hattı", "Mongo raw/live kayıt sayısı ve iki consumer group lag=0 kanıtını göster.")]
    for i, (n, h, b) in enumerate(demo):
        x = 0.7 + (i % 3) * 4.15; y = 1.45 + (i // 3) * 2.25
        card(s, f"{n}. {h}", b, x, y, 3.7, 1.65, BLUE if i < 3 else GREEN)
    text(s, "Demo sırasında hata olursa: log göstermek yerine önce /health, ardından consumer lag ve servis loglarıyla problemi daralt.", 0.9, 6.3, 11.5, 0.35, 15, RED, True, PP_ALIGN.CENTER)

    # 22 lessons
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "17. Staj boyunca ne öğrendim?", "Tek tek araçlardan çok, araçların birlikte oluşturduğu sistem davranışını öğrendim.", 22)
    lessons = [
        ("Sistem tasarımı", "Sorumlulukları ayırmak; veri üretimi, depolama, API ve arayüzü birbirinden bağımsızlaştırmak."),
        ("Dağıtık veri", "Topic, consumer group, offset, lag, yeniden deneme ve idempotency kavramlarını gerçek akışta kullanmak."),
        ("Veri modelleme", "Ham geçmiş ile güncel durum için farklı koleksiyon, index ve TTL stratejileri tasarlamak."),
        ("Ürün geliştirme", "REST snapshot ile WebSocket’in birbirini tamamladığını; fallback ve performansın kullanıcı deneyimini belirlediğini görmek."),
        ("Operasyon", "Container, healthcheck, backup/restore, monitoring, release kanıtı ve güvenli geri alma süreçlerini uygulamak."),
        ("Analitik", "Bronze/Silver/Gold, veri kalitesi, incremental refresh ve Iceberg snapshot/time-travel yaklaşımını uygulamak."),
    ]
    for i, (h, b) in enumerate(lessons):
        x, y = 0.7 + (i % 2) * 6.05, 1.35 + (i // 2) * 1.65
        card(s, h, b, x, y, 5.5, 1.3, [BLUE, GREEN, PURPLE, ORANGE, RED, CYAN][i])

    # 23 next
    s = P.slides.add_slide(P.slide_layouts[6]); title(s, "18. Sıradaki teknik adımlar", "Mevcut çalışma, güvenli release kabulüyle tamamlanacak; ardından iyileştirmeler ölçülerek ele alınacak.", 23)
    card(s, "Öncelik 1 — RC.3 kabulü", "Temiz Kafka volume, offline/registry kurulumları ve çok mimari güven zinciri kanıtlarını tamamlamak.", 0.8, 1.45, 3.65, 2.0, RED)
    card(s, "Öncelik 2 — operasyon", "DLQ izleme, metrik alarm eşiği ve release kabul matrisini gerçek hedeflerde doldurmak.", 4.85, 1.45, 3.65, 2.0, ORANGE)
    card(s, "Öncelik 3 — production kararları", "TLS, kimlik doğrulama, uygun harita tile sağlayıcısı ve ölçekleme planını eğitim ortamından ayrı değerlendirmek.", 8.9, 1.45, 3.65, 2.0, BLUE)
    text(s, "Bilinçli sınır: Mevcut yapı localhost eğitim projesidir; public internet yayını veya production SLA iddiası yoktur.", 0.95, 4.45, 11.45, 0.5, 17, RED, True, PP_ALIGN.CENTER)
    bullets(s, ["Önce doğrulanmış, geri alınabilir ve gözlemlenebilir sistem.", "Sonra performans / ölçekleme / yeni özellikler."], 3.05, 5.45, 7.3, 0.8, 18, NAVY)

    # 24 close
    s = P.slides.add_slide(P.slide_layouts[6])
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    rect(s, 0, 0, 13.333, 0.18, CYAN)
    text(s, "Sonuç", 0.78, 1.1, 3.0, 0.6, 30, WHITE, True)
    text(s, "Canlı uçuş verisini yalnızca haritada göstermedik;\nverinin güvenilir, izlenebilir ve analiz edilebilir bir yaşam döngüsünü kurduk.", 0.8, 2.0, 11.4, 1.1, 24, "D9EAF7", True)
    for i, (h, b, col) in enumerate([
        ("Canlı sistem", "OpenSky → Kafka → MongoDB → FastAPI → React", CYAN),
        ("Dayanıklılık", "Retry, DLQ, idempotent upsert, backup/restore", ORANGE),
        ("Analitik", "Spark, Bronze/Silver/Gold, Iceberg snapshot", PURPLE),
    ]):
        x = 0.85 + i * 4.1
        rect(s, x, 4.4, 3.5, 1.15, "1F4A6D", True, "2E648E")
        text(s, h, x + 0.2, 4.62, 3.1, 0.25, 15, col, True, PP_ALIGN.CENTER)
        text(s, b, x + 0.18, 4.97, 3.15, 0.28, 10, WHITE, False, PP_ALIGN.CENTER)
    text(s, "Teşekkürler • Sorularınız?", 0.8, 6.45, 11.7, 0.35, 18, WHITE, True, PP_ALIGN.CENTER)


add_content()
OUT.parent.mkdir(parents=True, exist_ok=True)
P.save(OUT)
print(f"Oluşturuldu: {OUT} ({len(P.slides)} slayt)")
